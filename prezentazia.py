import os
import sys
from pptx import Presentation
from pptx.util import Inches

def create_presentation():
    # Получаем текущую директорию программы
    current_dir = os.path.dirname(os.path.abspath(__file__))
    input_folder = current_dir
    output_pptx = os.path.join(current_dir, "presentation.pptx")

    # Создаем новую презентацию
    prs = Presentation()

    # Получаем список всех файлов PNG в текущей директории
    image_files = [f for f in os.listdir(input_folder) if f.endswith('.png')]

    # Для каждого файла PNG создаем слайд и добавляем изображение
    for image_file in image_files:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Используем пустой макет для слайда
        img_path = os.path.join(input_folder, image_file)
        left = top = Inches(1)
        pic = slide.shapes.add_picture(img_path, left, top)

    # Сохраняем презентацию
    prs.save(output_pptx)

if __name__ == "__main__":
    create_presentation()
